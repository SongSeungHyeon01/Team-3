"""PPTX / DOCX / XLSX 파싱"""
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook

def parse_pptx(file_path: str) -> dict:
    prs, pages = Presentation(file_path), []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        if slide.has_notes_slide: texts.append(slide.notes_slide.notes_text_frame.text)
        pages.append({"page_no": i, "text": "\n".join(texts), "image_path": None})
    return {"pages": pages}

def parse_docx(file_path: str) -> dict:
    doc = Document(file_path)
    text = "\n".join(p.text for p in doc.paragraphs)
    return {"pages": [{"page_no": 1, "text": text, "image_path": None}]}

def parse_xlsx(file_path: str) -> dict:
    wb, pages = load_workbook(file_path, data_only=True), []
    for i, ws in enumerate(wb.worksheets, 1):
        rows = [[str(c.value or "") for c in row] for row in ws.iter_rows()]
        text = "\n".join("\t".join(r) for r in rows)
        pages.append({"page_no": i, "text": text, "image_path": None})
    return {"pages": pages}
